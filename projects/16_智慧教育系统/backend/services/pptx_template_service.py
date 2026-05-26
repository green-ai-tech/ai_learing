"""PPTX 模板组件服务。"""

from __future__ import annotations

from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class PptxTemplateService:
    """封装课件模板化 slide 结构。"""

    navy = RGBColor(31, 41, 55)
    blue = RGBColor(37, 99, 235)
    green = RGBColor(20, 184, 166)
    gray = RGBColor(107, 114, 128)
    light = RGBColor(243, 244, 246)
    white = RGBColor(255, 255, 255)

    def __init__(self, presentation: Presentation):
        self.prs = presentation
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_cover(self, title: str, subtitle: str) -> None:
        slide = self._blank_slide()
        self._background_band(slide, self.navy)
        self._add_text(slide, title, Inches(0.9), Inches(1.8), Inches(9.5), Inches(0.9), 38, self.white, bold=True)
        self._add_text(slide, subtitle, Inches(0.95), Inches(2.85), Inches(8.5), Inches(0.4), 17, RGBColor(209, 213, 219))
        self._add_accent_block(slide, Inches(10.2), Inches(1.2), Inches(2.2), Inches(4.8), self.blue)

    def add_intro(self, title: str, description: str, bullets: Iterable[str]) -> None:
        slide = self._content_slide(title)
        self._add_text(slide, description, Inches(0.8), Inches(1.35), Inches(5.4), Inches(1.6), 18, self.navy)
        self._bullet_panel(slide, "课程画像", bullets, Inches(6.8), Inches(1.25), Inches(5.4), Inches(4.8), self.blue)

    def add_agenda(self, title: str, chapters: list[dict]) -> None:
        slide = self._content_slide(title)
        for index, chapter in enumerate(chapters[:8], 1):
            y = Inches(1.15 + (index - 1) * 0.62)
            self._number_badge(slide, index, Inches(0.9), y)
            self._add_text(slide, chapter.get("title", ""), Inches(1.55), y, Inches(8.8), Inches(0.42), 17, self.navy, bold=True)
            self._add_text(slide, f"{chapter.get('hours', '')}课时", Inches(10.55), y, Inches(1.4), Inches(0.35), 12, self.gray)

    def add_requirements(self, title: str, requirements: list[str], methods: list[str], assessments: list[str]) -> None:
        slide = self._content_slide(title)
        self._bullet_panel(slide, "教学要求", requirements, Inches(0.75), Inches(1.25), Inches(3.8), Inches(4.9), self.blue)
        self._bullet_panel(slide, "教学方法", methods, Inches(4.75), Inches(1.25), Inches(3.8), Inches(4.9), self.green)
        self._bullet_panel(slide, "考核方式", assessments, Inches(8.75), Inches(1.25), Inches(3.8), Inches(4.9), RGBColor(124, 58, 237))

    def add_chapter(self, chapter: dict) -> None:
        slide = self._blank_slide()
        self._background_band(slide, self.blue)
        self._add_text(slide, chapter.get("title", ""), Inches(0.9), Inches(2.1), Inches(8.8), Inches(0.7), 32, self.white, bold=True)
        self._add_text(slide, chapter.get("description", ""), Inches(0.95), Inches(3.0), Inches(8.2), Inches(1.0), 17, RGBColor(219, 234, 254))
        self._add_text(slide, f"{chapter.get('hours', '')} 课时", Inches(10.2), Inches(5.7), Inches(1.7), Inches(0.5), 22, self.white, bold=True)

    def add_theory(self, chapter: dict) -> None:
        slide = self._content_slide(f"{chapter.get('title', '')} 理论讲述")
        goals = chapter.get("teaching_goals") or []
        sections = [section.get("title", "") for section in chapter.get("sections", [])]
        self._bullet_panel(slide, "学习目标", goals, Inches(0.75), Inches(1.25), Inches(5.5), Inches(4.9), self.blue)
        self._bullet_panel(slide, "知识结构", sections, Inches(6.75), Inches(1.25), Inches(5.5), Inches(4.9), self.green)

    def add_explanation(self, chapter: dict) -> None:
        slide = self._content_slide(f"{chapter.get('title', '')} 讲解设计")
        section = (chapter.get("sections") or [{}])[0]
        self._add_text(slide, section.get("description", chapter.get("description", "")), Inches(0.85), Inches(1.25), Inches(6), Inches(1.5), 18, self.navy)
        points = []
        for point in section.get("knowledge_points", []):
            points.append(point.get("description") or point.get("name", ""))
        self._bullet_panel(slide, "讲解要点", points[:5], Inches(7.1), Inches(1.15), Inches(4.9), Inches(4.8), self.blue)
        self._process_bar(slide, ["导入", "讲解", "练习", "反馈"], Inches(0.85), Inches(4.85))

    def add_mixed(self, chapter: dict) -> None:
        slide = self._content_slide(f"{chapter.get('title', '')} 图文混编")
        self._image_placeholder(slide, Inches(0.8), Inches(1.35), Inches(5.2), Inches(4.15))
        section_titles = [section.get("title", "") for section in chapter.get("sections", [])]
        self._bullet_panel(slide, "课堂活动", section_titles, Inches(6.45), Inches(1.35), Inches(5.4), Inches(4.15), self.green)

    def _blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _content_slide(self, title: str):
        slide = self._blank_slide()
        self._add_text(slide, title, Inches(0.65), Inches(0.35), Inches(11.8), Inches(0.45), 24, self.navy, bold=True)
        self._add_line(slide, Inches(0.65), Inches(0.92), Inches(12.0))
        return slide

    def _background_band(self, slide, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(1, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _add_accent_block(self, slide, left, top, width, height, color: RGBColor) -> None:
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    def _add_text(self, slide, text: str, left, top, width, height, size: int, color: RGBColor, bold: bool = False):
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = str(text or "")
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        paragraph.font.name = "Microsoft YaHei"
        return box

    def _bullet_panel(self, slide, title: str, items: Iterable[str], left, top, width, height, color: RGBColor) -> None:
        panel = slide.shapes.add_shape(1, left, top, width, height)
        panel.fill.solid()
        panel.fill.fore_color.rgb = self.light
        panel.line.color.rgb = RGBColor(229, 231, 235)
        self._add_text(slide, title, left + Inches(0.25), top + Inches(0.18), width - Inches(0.5), Inches(0.35), 16, color, bold=True)
        box = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.75), width - Inches(0.6), height - Inches(0.95))
        frame = box.text_frame
        frame.clear()
        for index, item in enumerate([item for item in items if item][:6]):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(item)
            paragraph.level = 0
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.navy
            paragraph.font.name = "Microsoft YaHei"

    def _number_badge(self, slide, number: int, left, top) -> None:
        shape = slide.shapes.add_shape(9, left, top, Inches(0.38), Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.blue
        shape.line.fill.background()
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.text = str(number)
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = self.white

    def _add_line(self, slide, left, top, width) -> None:
        line = slide.shapes.add_shape(1, left, top, width, Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(229, 231, 235)
        line.line.fill.background()

    def _process_bar(self, slide, labels: list[str], left, top) -> None:
        for index, label in enumerate(labels):
            x = left + Inches(index * 1.55)
            self._number_badge(slide, index + 1, x, top)
            self._add_text(slide, label, x + Inches(0.48), top + Inches(0.02), Inches(0.9), Inches(0.3), 13, self.navy, bold=True)

    def _image_placeholder(self, slide, left, top, width, height) -> None:
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(219, 234, 254)
        shape.line.color.rgb = RGBColor(147, 197, 253)
        self._add_text(slide, "图示 / 案例 / 流程", left + Inches(1.35), top + Inches(1.8), Inches(2.6), Inches(0.45), 18, self.blue, bold=True)

